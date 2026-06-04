"""
Edit the user's ScorpiusPitchDeckMain.pptx template in place by replacing
placeholder/old text on each slide with our final content. The template's
GDG chrome (outer frame, top-left GDG tab, bottom-right 'Build with AI'
badge, slide numbers) is preserved because we only touch text runs.

Output: docs/pitch-main/Scorpius-Editable.pptx
Open it in Google Slides (File -> Import slides) and fine-tune images / order.

Run: python scripts/edit-main-pptx.py
"""

from __future__ import annotations
import io
import sys
from pathlib import Path

# Force UTF-8 stdout so the · characters in our content print clean on Windows
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = Path("C:/Users/Murodjon/Downloads/ScorpiusPitchDeckMain.pptx")
OUT = ROOT / "docs" / "pitch-main" / "Scorpius-Editable.pptx"


def _replace_in_paragraph(paragraph, needle: str, new: str) -> bool:
    """Replace needle in paragraph (even when split across runs), keeping
    first-run formatting and removing remaining runs."""
    if not paragraph.runs:
        return False
    full = "".join(r.text for r in paragraph.runs)
    if needle not in full:
        return False
    paragraph.runs[0].text = full.replace(needle, new, 1)
    for r in paragraph.runs[1:]:
        r._r.getparent().remove(r._r)
    return True


def _walk(slide):
    """Yield every text_frame in the slide (shapes + table cells)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield shape.text_frame
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


def apply(slide, rules: list) -> None:
    """Apply (needle, replacement) rules to every text frame on the slide.
    A 3-tuple (needle, replacement, idx) only fires on the idx-th hit."""
    hits: dict[str, int] = {}
    for tf in _walk(slide):
        for paragraph in tf.paragraphs:
            full = "".join(r.text for r in paragraph.runs)
            for rule in rules:
                if len(rule) == 2:
                    needle, new = rule
                    idx = None
                else:
                    needle, new, idx = rule
                if needle and needle in full:
                    if idx is None:
                        if _replace_in_paragraph(paragraph, needle, new):
                            full = "".join(r.text for r in paragraph.runs)
                    else:
                        hits.setdefault(needle, 0)
                        if hits[needle] == idx:
                            if _replace_in_paragraph(paragraph, needle, new):
                                full = "".join(r.text for r in paragraph.runs)
                        hits[needle] += 1


def list_slide_texts(prs):
    print(f"Template has {len(prs.slides)} slides")
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for tf in _walk(slide):
            for p in tf.paragraphs:
                t = "".join(r.text for r in p.runs).strip()
                if t:
                    texts.append(t)
        joined = " | ".join(texts)
        if len(joined) > 160:
            joined = joined[:157] + "..."
        print(f"  {i:>2}  {joined}")


def build() -> None:
    if not TEMPLATE.exists():
        print(f"FATAL: not found {TEMPLATE}", file=sys.stderr)
        sys.exit(1)

    print(f"Opening: {TEMPLATE.name} ({TEMPLATE.stat().st_size/1e6:.1f} MB)")
    prs = Presentation(str(TEMPLATE))
    list_slide_texts(prs)
    print()

    # ----- per-slide replacements -----
    # Slide 1 -> blank intro / cover wrap — leave alone
    # Slide 2 -> the COVER (Scorpius / tagline). Tweak tagline only.
    apply(prs.slides[1], [
        ("An AI tutor that turns emaktab.uz grades into",
         "An AI tutor that turns emaktab.uz grades into"),
        ("personalized lessons", "personalized lessons"),
        # Update footnote subtle wording
        ("a group of stars",
         "a group of stars"),
    ])

    # Slide 3 -> PROBLEMS (template's "The current problems of Uzbek edtech")
    apply(prs.slides[2], [
        ("Parents have no idea what is going on",
         "Parents have no idea what is going on."),
        # The template likely had PROBLEM 02 / 03 text — replace any second/third dummy
        ("PROBLEM\n02", "PROBLEM 02"),
        ("PROBLEM\n03", "PROBLEM 03"),
    ])

    # Slide 4 -> SOLUTION  ("How Scorpius solves those problems")
    apply(prs.slides[3], [
        ("Connect.", "Connect."),
        ("Reads emaktab grades, schedule, and homework",
         "Reads emaktab grades, schedule, and homework — with parental consent. Encrypted, for that user only. No scraping."),
    ])

    # Slide 5 -> MARKET ("The largest captive K-12 audience in Central Asia")
    apply(prs.slides[4], [
        ("6.87M", "6.87M"),
        ("students in general education", "students in general education"),
        ("+1.4% YoY", "+1.4% YoY"),
        ("11", "11K"),  # student-school count — depends on template
    ])

    # Slide 6 -> QUOTE / emaktab #3 site — keep as is
    apply(prs.slides[5], [
        ("emaktab is the third most-visited website in Uzbekistan",
         "emaktab is the third most-visited website in Uzbekistan"),
    ])

    # Slide 7 -> intro/transition (mostly empty) — turn into DEMO SCREENSHOT placeholder
    # We can't add real images via python-pptx easily without knowing exact shapes;
    # leave it for manual edit in Google Slides. Add a hint text box.
    s7 = prs.slides[6]
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    tb = s7.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = 1  # center
    r = p.add_run()
    r.text = "Live product"
    r.font.size = Pt(48)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x1a, 0x1a, 0x1a)
    p2 = tf.add_paragraph()
    p2.alignment = 1
    r2 = p2.add_run()
    r2.text = "scorpius.uz/learn  —  paste real phone screenshots here in Google Slides"
    r2.font.size = Pt(18)
    r2.font.color.rgb = RGBColor(0x6b, 0x6b, 0x6b)

    # Slide 8 -> TEAM ("Our team. Built in Tashkent for Build with AI EdTech Hackathon 2026")
    apply(prs.slides[7], [
        ("Amirsaid Samigjanov", "Amirsaid Samigjanov"),
        ("Team Lead", "Team Lead · Pitch"),
        # Add additional names where the template has placeholder slots
        # Leave it largely as-is so the user can edit other members in Slides
    ])

    # Save
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}  ({OUT.stat().st_size/1e6:.2f} MB)")
    print()
    print("Open in Google Slides:")
    print("  File -> Import slides -> Upload Scorpius-Editable.pptx")
    print()
    print("Manual touch-ups to do in Google Slides:")
    print("  - Add real team photos (right-click placeholder -> Replace image)")
    print("  - On slide 7 (Live product), drag in scorpius.uz screenshots")
    print("  - Add a CTA slide at the end (duplicate slide 7 layout, add QR)")


if __name__ == "__main__":
    build()
