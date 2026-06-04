"""
Build the Scorpius pitch deck from the official GDG BWAI 2026 template.

Pipeline:
  1. Open the 108-slide template
  2. Apply text replacements ONLY on slides we plan to keep (so the wrong slide
     doesn't accidentally absorb our content)
  3. Trim down to the kept slides
  4. Reorder them into our pitch sequence (1..10)
  5. Append an 11th slide for the team (4 photo + bio cells)
  6. Save to docs/Scorpius-Pitch-Deck.pptx

Replacement is multi-run aware: when a needle is split across runs in a
paragraph, the whole paragraph text is rewritten while preserving the first
run's formatting.

Run: python scripts/build-pitch-pptx.py
"""

from __future__ import annotations
import io
import sys
from pathlib import Path

# Windows console defaults to cp1252; force UTF-8 so · and other chars print fine
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "BWAI 2026 Master Deck - Template.pptx"
OUT = ROOT / "docs" / "Scorpius-Pitch-Deck.pptx"
TEAM_DIR = ROOT / "public" / "team"

# (template_slide_number_1based, pitch_position_1based, label)
PITCH = [
    (2,   1,  "01 . Cover"),
    (3,   2,  "02 . Cloud Credit (mandatory)"),
    (4,   3,  "03 . Big Idea"),
    (19,  4,  "04 . Parent quote"),
    (53,  5,  "05 . Market (2.4M / 0)"),
    (48,  6,  "06 . Validated demand"),
    (81,  7,  "07 . What they do today (chart)"),
    (103, 8,  "08 . Live demo screenshots"),
    (70,  9,  "09 . What's ours (bullets)"),
    (75,  10, "10 . CTA + QR + video"),
]

# Per slide replacements. Two forms supported:
#   ("needle", "replacement")  - global; replaces every paragraph match (so guard against duplicates)
#   ("needle", "replacement", n)  - tuple of len 3: replace only the nth (0-based) paragraph match
REPLACEMENTS: dict[int, list] = {
    # 1. Cover
    2: [
        ("Topic:", "AI Tutor for emaktab.uz"),
        ("Your team name", "Scorpius — Murodjon . Amirsaid . Zaynobiddin"),
    ],
    # 2. Cloud Credit — leave alone
    3: [],
    # 3. Big Idea — replace the whole stacked headline
    4: [
        ("The", "emaktab"),
        ("Big", "tells you"),
        ("Idea", "the grade."),
        ("Goes Here.", "We teach you why."),
        (
            "This is where the subtitle can be displayed.",
            "Scorpius — an AI tutor that turns emaktab.uz grades into personal lessons. Socratic, in Uzbek, 8 minutes per topic.",
        ),
    ],
    # 4. Parent quote — replace the stacked Dieter Rams quote
    19: [
        ("Good design", "emaktab tells me"),
        ("is as", "my child got 2 today."),
        ("little", "It doesn't tell me"),
        ("design", "why,"),
        ("as possible.", "or what to do."),
        ("- Dieter Rams", "— Parent, Tashkent . cust-dev May 2026"),
    ],
    # 5. Market 2.4M / 0 — template slide 53 has two stats side by side
    53: [
        ("86%", "2.4M", 0),  # first big number -> 2.4M
        ("Statistic caption", "students on emaktab daily", 0),
        ("Statistic caption", "  ", 1),
        ("Learnings", "The market is locked in. The product layer is empty."),
        (
            "This is body copy. Bringing developers together in-person and online. Stay in the know about upcoming events, catch up on content from past events, and meet other developers in the community.",
            "2.4 million students log into emaktab every day. Zero get a personalised lesson for the topic they just missed. That layer above the gradebook is the gap Scorpius fills.",
        ),
    ],
    # 6. Validated demand — template slide 48 has one big % stat + learnings body
    48: [
        ("97%", "100%"),
        (
            "Statistic caption this is body copy and it goes a little like this",
            "of 70 students said YES to an AI tutor today. Zero rejection.",
        ),
        ("Learnings", "Validated this morning"),
        (
            "This is body copy. Bringing developers together in-person and online. Stay in the know about upcoming events, catch up on content from past events, and meet other developers in the community.",
            "70 cust-dev responses in 4 hours via Telegram. 44% would use daily, 29% sometimes, 27% want to try. Half (35) left their Telegram username for the beta. Source: forms.gle/zE9bryCvewNKL5jUA",
        ),
    ],
    # 7. Chart — slide 81 has 8 "Short label" texts; we replace per-index below.
    81: [
        ("Simple chart", "When students don't understand homework, they…"),
    ],
    # 8. iPhone device frame — slide 103 has 2 image placeholders
    103: [
        (
            "Right click and “replace image”",
            "[ paste screenshot of scorpius.uz ]",
        ),
        ("DEVICE FRAME - IPHONE 14", "Live now on scorpius.uz"),
    ],
    # 9. Bullets
    70: [
        ("Left Aligned Title", "Not a model wrapper. The IP is ours."),
        (
            "Bullet one",
            "Card DSL — 12 typed lesson card variants (mcq, discover, sequence, diagram, ask, story…)",
        ),
        (
            "Bullet two, adjust size to match length of content",
            "Model adapter — swap OpenAI ↔ Gemini with one env var",
        ),
        (
            "Bullet three",
            "Uzbek curriculum RAG — Grade-6 Maths + Physics from real textbooks via gpt-5.1 vision",
        ),
        (
            "Bullet four, a little longer example",
            "SHA-keyed Firestore cache — one student pays, every student after is free",
        ),
        (
            "A final",
            "Govt support — Fergana viloyati Maktabgacha va Maktab ta'lim bo'limi backs us",
        ),
    ],
    # 10. CTA + QR + video
    75: [
        (
            "Simple quote or statement goes here. Ideally limit to four or five lines max.",
            "Scan the QR. Try it. Tell us. Birinchi 100 ga umrbod 50%. scorpius.uz/gdg",
        ),
    ],
}

# Chart labels for slide 81. 8 in order — replaces every "Short label" / "Short Label" left to right.
CHART_LABELS = [
    "ChatGPT / AI  50%",
    "YouTube  20%",
    "Classmate  11%",
    "Other / give up  9%",
    "Family  3%",
    "Tutor  3%",
    "Khan Academy  3%",
    "Nobody  1%",
]


def _replace_in_paragraph(paragraph, needle: str, new: str) -> bool:
    """Replace 'needle' in a paragraph, even when split across runs.
    Returns True if a replacement was made.
    Preserves the first run's formatting; collapses other runs."""
    if not paragraph.runs:
        return False
    full = "".join(r.text for r in paragraph.runs)
    if needle not in full:
        return False
    new_full = full.replace(needle, new, 1)
    # Keep first run, clear others
    first = paragraph.runs[0]
    first.text = new_full
    # Remove subsequent runs to avoid duplicate text appended
    for r in paragraph.runs[1:]:
        r._r.getparent().remove(r._r)
    return True


def _apply_rules_to_text_frame(tf, rules: list) -> None:
    """Apply replacement rules to a text frame, paragraph by paragraph."""
    # Per-needle paragraph hit counter for indexed rules
    hit_counts: dict[str, int] = {}
    for paragraph in tf.paragraphs:
        for rule in rules:
            if len(rule) == 2:
                needle, new = rule
                idx = None
            else:
                needle, new, idx = rule
            if idx is not None:
                hit_counts.setdefault(needle, 0)
                # Look first if this paragraph contains the needle
                full = "".join(r.text for r in paragraph.runs)
                if needle in full:
                    if hit_counts[needle] == idx:
                        _replace_in_paragraph(paragraph, needle, new)
                    hit_counts[needle] += 1
            else:
                _replace_in_paragraph(paragraph, needle, new)


def _walk_text_frames(slide):
    """Yield every text frame on the slide (top-level shapes + table cells)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield shape.text_frame
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


def _apply_chart_labels(slide):
    """Slide 81 — replace 'Short label' / 'Short Label' in document order with CHART_LABELS."""
    idx = 0
    for tf in _walk_text_frames(slide):
        for paragraph in tf.paragraphs:
            for label_needle in ("Short label", "Short Label"):
                full = "".join(r.text for r in paragraph.runs)
                if label_needle in full and idx < len(CHART_LABELS):
                    _replace_in_paragraph(paragraph, label_needle, CHART_LABELS[idx])
                    idx += 1


def _drop_slides(prs: Presentation, keep_0based: set[int]) -> None:
    """Remove slides whose 0-based index is not in keep_0based, then return the
    presentation. Operates on the XML sldIdLst."""
    slides_el = prs.slides._sldIdLst
    id_list = list(slides_el)
    for i, sld_id in enumerate(id_list):
        if i not in keep_0based:
            rId = sld_id.get(qn("r:id"))
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
            slides_el.remove(sld_id)


def _reorder_slides(prs: Presentation, new_order_old_indices: list[int]) -> None:
    """Reorder slides. `new_order_old_indices[i]` = the original 0-based index
    that should land at position i in the new order."""
    slides_el = prs.slides._sldIdLst
    current = list(slides_el)
    for el in current:
        slides_el.remove(el)
    for old_idx in new_order_old_indices:
        slides_el.append(current[old_idx])


def build() -> None:
    if not TEMPLATE.exists():
        print(f"FATAL: template not found at {TEMPLATE}", file=sys.stderr)
        sys.exit(1)

    print(f"Opening template: {TEMPLATE.name} ({TEMPLATE.stat().st_size / 1e6:.1f} MB)")
    prs = Presentation(str(TEMPLATE))
    print(f"Template has {len(prs.slides)} slides")

    # 1. Apply text replacements on the kept slides BEFORE we touch slide order/count
    keep_set_0based = {tpl - 1 for tpl, _pos, _label in PITCH}
    for tpl, pos, label in PITCH:
        slide = prs.slides[tpl - 1]
        print(f"  {label} (template slide {tpl}) -> pitch position {pos}")
        rules = REPLACEMENTS.get(tpl, [])
        for tf in _walk_text_frames(slide):
            _apply_rules_to_text_frame(tf, rules)
        if tpl == 81:
            _apply_chart_labels(slide)

    # 2. Drop everything not in our keep set
    print("Trimming slides…")
    _drop_slides(prs, keep_set_0based)
    print(f"  -> {len(prs.slides)} slides remain")

    # 3. Reorder. After dropping, the kept slides are in template order. We need
    #    them in pitch order. Build the mapping:
    #    after dropping, slides are at indices 0..N-1 in template-sorted-by-original-index.
    sorted_pitch_by_tpl = sorted(PITCH, key=lambda x: x[0])  # template order after drop
    # current_idx_of_tpl[tpl] = 0..N-1 = its position after the drop
    current_idx_of_tpl = {tpl: i for i, (tpl, _pos, _lbl) in enumerate(sorted_pitch_by_tpl)}
    pitch_sorted = sorted(PITCH, key=lambda x: x[1])  # by pitch position
    new_order = [current_idx_of_tpl[tpl] for tpl, _pos, _lbl in pitch_sorted]
    print(f"Reorder mapping (new position -> old index after trim): {new_order}")
    _reorder_slides(prs, new_order)

    # Save (without team slide yet) — team gets appended in a separate pass
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved trimmed + ordered deck: {OUT}")

    # 4. Append team slide
    _add_team_slide(OUT)
    print("Team slide appended (11 slides total)")
    print()
    print("DONE")
    print(f"Output: {OUT}")
    print(f"Size: {OUT.stat().st_size / 1e6:.1f} MB")


# ----- Team slide -----------------------------------------------------------

# Ladies first per Boss's instruction. Replace these names/bios after the user
# confirms — placeholders are clearly labelled so it's obvious what to edit.
TEAM = [
    {
        "name": "[Add name]",
        "role": "Product . Operations",
        "bio": "Replace this with her real bio. Where she studies / works . what she has shipped . why she is on this team. 2-3 sentences max.",
        "photo": "member1.jpg",
    },
    {
        "name": "Murodjon Tolipov",
        "role": "Founder . Build",
        "bio": "Founder. Built the Scorpius platform end-to-end this hackathon (Next.js 16, Firebase, OpenAI). Previously: Humopedia — an Uzbek-first knowledge platform.",
        "photo": "murodjon.jpg",
    },
    {
        "name": "Amirsaid",
        "role": "Team lead . Pitch",
        "bio": "Team lead and pitcher. Drives narrative, product positioning, and customer research. Shipped the 70-response cust-dev survey this morning.",
        "photo": "amirsaid.jpg",
    },
    {
        "name": "Zaynobiddin Shakhabiddinov",
        "role": "CTO . Infrastructure",
        "bio": "CTO and infra owner. Shipped Grade-6 physics curriculum and the learning-tree feature this hackathon. Da Vinci Resolve editor — owns the demo video.",
        "photo": "zaynobiddin.jpg",
    },
]


def _add_team_slide(pptx_path: Path) -> None:
    prs = Presentation(str(pptx_path))

    # Pick the blankest layout
    layout = None
    for lay in prs.slide_layouts:
        if len(lay.placeholders) <= 1:
            layout = lay
            break
    if layout is None:
        layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)

    sw, sh = prs.slide_width, prs.slide_height
    margin = Inches(0.5)

    # Title
    title_box = slide.shapes.add_textbox(margin, Inches(0.4), sw - 2 * margin, Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "The team"
    r.font.name = "Fraunces"
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x21, 0x21, 0x21)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(margin, Inches(1.25), sw - 2 * margin, Inches(0.5))
    p = subtitle_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Scorpius . Tashkent . GDG Build with AI 2026"
    r.font.name = "Google Sans"
    r.font.size = Pt(14)
    r.font.color.rgb = RGBColor(0x59, 0x59, 0x59)

    # 2×2 grid
    grid_top = Inches(2.0)
    grid_h = sh - grid_top - Inches(0.5)
    grid_w = sw - 2 * margin
    cols, rows = 2, 2
    cell_w = grid_w / cols
    cell_h = grid_h / rows

    photo_w = Inches(1.7)
    photo_h = Inches(1.7)
    text_w = cell_w - photo_w - Inches(0.5)

    from pptx.enum.shapes import MSO_SHAPE

    for i, m in enumerate(TEAM):
        col = i % cols
        row = i // cols
        cx = margin + col * cell_w
        cy = grid_top + row * cell_h

        # Photo or placeholder
        photo_path = TEAM_DIR / m["photo"]
        if photo_path.exists():
            slide.shapes.add_picture(str(photo_path), cx, cy, photo_w, photo_h)
        else:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, photo_w, photo_h
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
            box.line.color.rgb = RGBColor(0xD4, 0xCD, 0xB9)
            box.line.width = Pt(0.75)
            ftf = box.text_frame
            ftf.word_wrap = True
            fp = ftf.paragraphs[0]
            fp.alignment = PP_ALIGN.CENTER
            fr = fp.add_run()
            fr.text = f"Drop\n{m['photo']}\ninto public/team/"
            fr.font.name = "Google Sans"
            fr.font.size = Pt(8.5)
            fr.font.color.rgb = RGBColor(0x9B, 0x95, 0x8A)

        # Text
        text_left = cx + photo_w + Inches(0.25)
        text_top = cy
        tbox = slide.shapes.add_textbox(text_left, text_top, text_w, photo_h)
        ttf = tbox.text_frame
        ttf.word_wrap = True

        p1 = ttf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run()
        r1.text = m["name"]
        r1.font.name = "Fraunces"
        r1.font.bold = True
        r1.font.size = Pt(18)
        r1.font.color.rgb = RGBColor(0x21, 0x21, 0x21)

        p2 = ttf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = m["role"]
        r2.font.name = "Google Sans"
        r2.font.size = Pt(10)
        r2.font.color.rgb = RGBColor(0x42, 0x85, 0xF4)  # Google Blue
        r2.font.bold = True

        p3 = ttf.add_paragraph()
        p3.alignment = PP_ALIGN.LEFT
        p3.space_before = Pt(6)
        r3 = p3.add_run()
        r3.text = m["bio"]
        r3.font.name = "Google Sans"
        r3.font.size = Pt(9.5)
        r3.font.color.rgb = RGBColor(0x59, 0x59, 0x59)

    prs.save(str(pptx_path))


if __name__ == "__main__":
    build()
